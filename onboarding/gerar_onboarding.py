from pptx import Presentation
from copy import deepcopy

# Dados do colaborador
dados = {
    "{{NOME}}": "João da Silva",
    "{{USUARIO}}": "joao.silva",
    "{{SENHA}}": "Senha@312!",
}

# Acessos do colaborador
acessos = {
    "{{ACESSO_REDE}}": True,
    "{{ACESSO_EMAIL}}": True,
    "{{ACESSO_FLUIG}}": True,
    "{{ACESSO_TOTVS}}": True,
    "{{ACESSO_EDATA}}": True,
    "{{ACESSO_SAG}}": True,
    "{{ACESSO_BI}}": True
}

# Pegando o arquivo de template
prs = Presentation("template_onboarding.pptx")

def substituir_texto(slide, substituicoes):
    for shape in slide.shapes: # Percorre cada elemento do slide
        if shape.has_text_frame: # Verificando se o elemento tem texto
            for para in shape.text_frame.paragraphs: # Verificando cada paragrafo
                for run in para.runs: # Verficando cada techo do texto
                    for marcador, valor in substituicoes.items():
                        if marcador in run.text:
                            run.text = run.text.replace(marcador,valor)

# Fazendo ele rodar pelos slides onde tem informações para trocar
# for i in [0, 3, 4, 5]:
#     substituir_texto(prs.slides[i], dados)

for i in [3, 4, 5]:
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for marcador, tem_acesso in acessos.items():
                        if marcador in run.text:
                            # Se tem acesso: ✅  |  Se não tem: ❌
                            run.text = run.text.replace(
                                marcador, "✅" if tem_acesso else ""
                            )

nome_arquivo = f"Onboarding - {dados['{{NOME}}']}.pptx" #criando o nome do arquivo
prs.save(nome_arquivo) # Salvando o arquivo
print(f"Arquivo gerado: {nome_arquivo}")