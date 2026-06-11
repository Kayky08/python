from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import comtypes.client
from pathlib import Path
import copy
import sys
import os
import io
import shutil

# ─── DADOS DO COLABORADOR ─────────────────────────────────────────────────────
# print("=== Gerador de Onboarding ===\n")

# def perguntar(acesso):
#     return input(f"  {acesso}? (s/n): ").strip().lower() == "s"

# ACESSOS = {
#     "Rede":      perguntar("Rede"),
#     "Office365": perguntar("Office365"),
#     "Fluig":     perguntar("Fluig"),
#     "Protheus":  perguntar("Protheus"),
#     "Edata":     perguntar("Edata"),
#     "SAG":       perguntar("SAG"),
#     "PowerBI":   perguntar("Power BI"),
#     "WMW":       perguntar("WMW"),
# }

# print("\nDefina as informações:\n")

# NOME          = input("Nome completo: ")
# USUARIO       = input("Usuário: ")
# SENHA         = input("Senha: ")
# USUARIO_EDATA = input("Usuário EDATA: ") if ACESSOS["Edata"] else ""
# USUARIO_SAG   = input("Usuário SAG: ")   if ACESSOS["SAG"]   else ""
# USUARIO_WMW   = input("Usuário WMW: ")   if ACESSOS["WMW"]   else ""

#---------------Teste----------------- 
ACESSOS = {
    "Rede":      "s",
    "Office365": "s",
    "Fluig":     "s",
    "Protheus":  "s",
    "Edata":     "s",
    "SAG":       "n",
    "PowerBI":   "n",
    "WMW":       "n"
}

NOME    = "Emelly Nachbal"
USUARIO = "emelly.nachbal"
SENHA   = "Qual@312!"
USUARIO_EDATA = "nemelly"
USUARIO_SAG = "emellynachbal"
USUARIO_WMW   = "000842"
#---------------Teste----------------- 

# ─── DEFINIÇÃO DE CADA ACESSO ─────────────────────────────────────────────────
DEFINICAO_ACESSOS = {
    "Rede": {
        "slide_origem": 3,
        "texto_shape":  "ct_rede",
        "imagens": [
            {"nome": "img_rede", "offset": -482429, "left": 9704070, "w": 1338263, "h": 1765194},
        ],
    },
    "Office365": {
        "slide_origem": 3,
        "texto_shape":  "ct_office",
        "imagens": [
            {"nome": "img_office", "offset": 537266, "left": 9924765, "w": 896873, "h": 807276},
        ],
    },
    "Fluig": {
        "slide_origem": 3,
        "texto_shape":  "ct_fluig",
        "imagens": [
            {"nome": "img_fluig", "offset": 486329, "left": 9704070, "w": 1098803, "h": 516635},
        ],
    },
    "Protheus": {
        "slide_origem": 4,
        "texto_shape":  "ct_protheus",
        "imagens": [
            {"nome": "img_protheus1", "offset": -395287, "left": 9412930,  "w": 714375, "h": 752475},
            {"nome": "img_protheus2", "offset":  428252, "left": 9972893,  "w": 771525, "h": 742950},
            {"nome": "img_protheus3", "offset": -357187, "left": 10550541, "w": 752475, "h": 714375},
        ],
    },
    "Edata": {
        "slide_origem": 4,
        "texto_shape":  "ct_edata",
        "imagens": [
            {"nome": "img_edata1", "offset":  835149, "left": 10045653, "w": 821837, "h": 769657},
            {"nome": "img_edata2", "offset":  -61635, "left": 9513846,  "w": 821837, "h": 731949},
            {"nome": "img_edata3", "offset":  -65836, "left": 10591686, "w": 711330, "h": 731948},
        ],
    },
    "SAG": {
        "slide_origem": 4,
        "texto_shape":  "ct_sag",
        "imagens": [
            {"nome": "img_sag", "offset": 451877, "left": 9940482, "w": 679896, "h": 945942},
        ],
    },
    "PowerBI": {
        "slide_origem": 5,
        "texto_shape":  "ct_bi",
        "imagens": [
            {"nome": "img_bi", "offset": -226036, "left": 9727692, "w": 1473708, "h": 1873896},
        ],
    },
    "WMW": {
        "slide_origem": 5,
        "texto_shape":  "ct_wmw",
        "imagens": [
            {"nome": "img_wmw", "offset": 500000, "left": 9978633, "w": 1021980, "h": 904059},
        ],
    },
}

CAPACIDADE_POR_SLIDE = 3
TOPO_INICIAL         = 1496968
ESPACO_ENTRE_BLOCOS  = 200000


# ─── FUNÇÕES ──────────────────────────────────────────────────────────────────

def substituir_texto(slide, usuario, senha, nome, usuario_edata="", usuario_sag="", usuario_wmw = ""):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = (run.text
                    .replace("{{USUARIO}}", usuario)
                    .replace("{{SENHA}}", senha)
                    .replace("{{NOME}}", nome)
                    .replace("{{USUARIO_EDATA}}", usuario_edata)
                    .replace("{{USUARIO_SAG}}", usuario_sag)
                    .replace("{{USUARIO_WMW}}", usuario_wmw))


def remover_slide(prs, indice):
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[indice])


def limpar_slide_de_acessos(slide):
    """Remove todos os shapes de acesso, mantendo só cabeçalho e barra."""
    shapes_fixos = {"object 2", "object 3"}
    sp_tree = slide.shapes._spTree
    para_remover = [
        shape._element for shape in slide.shapes
        if shape.name not in shapes_fixos
    ]
    for el in para_remover:
        sp_tree.remove(el)


def copiar_shape_com_imagem(shape_origem, slide_origem, slide_destino):
    """
    Copia um shape de imagem garantindo que o arquivo de imagem
    seja registrado corretamente no slide de destino.

    O problema original: copy.deepcopy copia o XML com o rId (ex: rId4),
    mas no slide de destino rId4 aponta para outra imagem.

    A solução: lemos os bytes reais da imagem pelo relationship do slide
    de origem, adicionamos esses bytes como um novo part no slide de destino
    e atualizamos o rId no XML copiado para apontar para o novo part.
    """
    # Namespace do atributo r:embed que guarda o rId da imagem no XML
    R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

    # 1. Encontra o elemento <a:blip> no XML do shape — é onde fica o rId
    blip = shape_origem._element.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    )
    if blip is None:
        # Não é uma imagem com blip (ex: shape de fundo) — copia normalmente
        slide_destino.shapes._spTree.append(copy.deepcopy(shape_origem._element))
        return

    rid_origem = blip.get(R_EMBED)  # ex: "rId4"

    # 2. Resolve o rId para o part (objeto que contém os bytes da imagem)
    rel_origem = slide_origem.part.rels.get(rid_origem)
    if rel_origem is None:
        slide_destino.shapes._spTree.append(copy.deepcopy(shape_origem._element))
        return

    image_part = rel_origem.target_part  # part com os bytes da imagem

    # 3. Adiciona a imagem no slide de destino e obtém um novo rId
    #    get_or_add_image_part registra o arquivo no pacote e cria o relationship
    #    no slide de destino, retornando o novo rId válido para esse slide
    img_bytes = io.BytesIO(image_part.blob)           # bytes reais da imagem
    _, novo_rid = slide_destino.part.get_or_add_image_part(img_bytes)

    # 4. Clona o XML do shape e atualiza o rId para o novo valor
    novo_elemento = copy.deepcopy(shape_origem._element)
    novo_blip = novo_elemento.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    )
    novo_blip.set(R_EMBED, novo_rid)  # substitui rId antigo pelo novo

    # 5. Insere o shape no slide de destino
    slide_destino.shapes._spTree.append(novo_elemento)


def adicionar_acesso_ao_slide(slide_destino, acesso_nome, novo_topo, prs_ref, definicao):
    """Copia texto + imagens do acesso para o slide de destino, já posicionados."""
    slide_origem = prs_ref.slides[definicao["slide_origem"]]
    shapes_origem = {s.name: s for s in slide_origem.shapes}

    # Copia o bloco de texto (não é imagem, deepcopy normal funciona)
    shape_texto = shapes_origem.get(definicao["texto_shape"])
    if shape_texto:
        novo_el = copy.deepcopy(shape_texto._element)
        slide_destino.shapes._spTree.append(novo_el)
        # Reposiciona
        for s in slide_destino.shapes:
            if s.name == definicao["texto_shape"]:
                s.top = novo_topo
                break

    # Copia cada imagem usando a função que preserva o vínculo correto
    for img_def in definicao["imagens"]:
        shape_img = shapes_origem.get(img_def["nome"])
        if not shape_img:
            continue

        copiar_shape_com_imagem(shape_img, slide_origem, slide_destino)

        # Reposiciona e redimensiona a imagem recém-copiada
        for s in slide_destino.shapes:
            if s.name == img_def["nome"]:
                s.top    = novo_topo + img_def["offset"]
                s.left   = img_def["left"]
                s.width  = img_def["w"]
                s.height = img_def["h"]
                break


# ─── LÓGICA PRINCIPAL ─────────────────────────────────────────────────────────

ativos = [nome for nome, tem in ACESSOS.items() if tem]

print(f"\nGerando onboarding para: {NOME}")
print(f"Acessos ativos ({len(ativos)}): {', '.join(ativos)}")
print("-" * 40)

if not ativos:
    print("Nenhum acesso selecionado.")
    exit()

grupos = []
i = 0
while i < len(ativos):
    grupos.append(ativos[i:i + CAPACIDADE_POR_SLIDE])
    i += CAPACIDADE_POR_SLIDE

print(f"Slides necessários: {len(grupos)}")
for g_idx, grupo in enumerate(grupos):
    print(f"  Slide {4 + g_idx}: {', '.join(grupo)}")

prs_ref = Presentation("template_onboarding.pptx")
prs     = Presentation("template_onboarding.pptx")

substituir_texto(prs.slides[0], USUARIO, SENHA, NOME, USUARIO_EDATA, USUARIO_SAG, USUARIO_WMW)
print("\n✓ Slide 1: nome preenchido")

slides_acesso_template = [3, 4, 5]

for g_idx, grupo in enumerate(grupos):
    slide_idx = slides_acesso_template[g_idx]
    slide = prs.slides[slide_idx]

    print(f"\nSlide {4 + g_idx} — {', '.join(grupo)}:")
    limpar_slide_de_acessos(slide)

    topo_atual = TOPO_INICIAL
    for acesso in grupo:
        definicao = DEFINICAO_ACESSOS[acesso]
        adicionar_acesso_ao_slide(slide, acesso, topo_atual, prs_ref, definicao)

        for shape in slide.shapes:
            if shape.name == definicao["texto_shape"]:
                topo_atual += shape.height + ESPACO_ENTRE_BLOCOS
                break

        print(f"  ✓ '{acesso}' posicionado")

    substituir_texto(slide, USUARIO, SENHA, NOME, USUARIO_EDATA, USUARIO_SAG, USUARIO_WMW)

slides_para_remover = slides_acesso_template[len(grupos):]
for idx in sorted(slides_para_remover, reverse=True):
    remover_slide(prs, idx)
    print(f"\n✗ Slide {idx + 1} removido")

nome_arquivo = f"Onboarding - {NOME}.pptx"
prs.save(nome_arquivo)

powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
powerpoint.Visible = 1

caminho_pptx = os.path.abspath(nome_arquivo)
caminho_pdf  = caminho_pptx.replace(".pptx", ".pdf")

deck = powerpoint.Presentations.Open(caminho_pptx)
deck.SaveAs(caminho_pdf, 32)
deck.Close()
powerpoint.Quit()

shutil.move(caminho_pdf,"ONBOARDING")
os.remove(caminho_pptx)

print(f"\n{'='*40}")
print(f"✅ Arquivo gerado: {nome_arquivo}")
print(f"✅ PDF gerado: {caminho_pdf}")
print(f"   Total de slides: {len(prs.slides)}")