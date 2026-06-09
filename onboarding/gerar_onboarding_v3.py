from pptx import Presentation

# ─── DADOS DO COLABORADOR ─────────────────────────────────────────────────────
print("=== Gerador de Onboarding ===\n")

NOME    = input("Nome completo: ")
USUARIO = input("Usuário: ")
SENHA   = input("Senha: ")

print("\nDefina os acessos (s = sim / n = não):\n")

def perguntar(acesso):
    resposta = input(f"  {acesso}? (s/n): ").strip().lower()
    return resposta == "s"

ACESSOS_SLIDE4 = {
    "Fluig":     perguntar("Fluig"),
    "Rede":      perguntar("Rede"),
    "Office365": perguntar("Office365"),
}

ACESSOS_SLIDE5 = {
    "Protheus": perguntar("Protheus"),
    "Edata":    perguntar("Edata"),
    "SAG":      perguntar("SAG"),
}

ACESSOS_SLIDE6 = {
    "PowerBI": perguntar("Power BI"),
}

# ─── MAPEAMENTO DE SHAPES ─────────────────────────────────────────────────────
# Cada acesso tem: shape de texto + lista de imagens associadas
# As imagens foram identificadas cruzando posições verticais no template

MAPA_ACESSOS = {
    # Slide 4
    "Fluig":     {"texto": "object 5",       "imagens": ["object 4"]},
    "Rede":      {"texto": "CaixaDeTexto 6", "imagens": []},
    "Office365": {"texto": "CaixaDeTexto 9", "imagens": ["object 6"]},
    # Slide 5
    "Protheus":  {"texto": "CaixaDeTexto 17","imagens": ["Imagem 16", "Imagem 9", "Imagem 13"]},
    "Edata":     {"texto": "CaixaDeTexto 18","imagens": ["Imagem 8", "Imagem 10"]},
    "SAG":       {"texto": "object 5",       "imagens": ["Imagem 12", "Imagem 14"]},
    # Slide 6
    "PowerBI":   {"texto": "object 5",       "imagens": ["Imagem 5"]},
}

# Ordem dos blocos em cada slide (de cima para baixo, pela posição no template)
ORDEM_SLIDE4 = ["Rede", "Office365", "Fluig"]
ORDEM_SLIDE5 = ["Protheus", "Edata", "SAG"]
ORDEM_SLIDE6 = ["PowerBI"]

# Posição do topo do primeiro bloco em cada slide (em EMU — unidade interna do pptx)
# Todos os blocos ativos vão se empilhar a partir desse ponto
TOPO_INICIAL_SLIDE4 = 1496968
TOPO_INICIAL_SLIDE5 = 1597636
TOPO_INICIAL_SLIDE6 = 1371600

# Espaçamento entre blocos (em EMU)
ESPACO_ENTRE_BLOCOS = 200000


# ─── FUNÇÕES ──────────────────────────────────────────────────────────────────

def substituir_texto(slide, usuario, senha, nome):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace("{{USUARIO}}", usuario)
                run.text = run.text.replace("{{SENHA}}", senha)
                run.text = run.text.replace("{{NOME}}", nome)


def remover_shape_por_nome(slide, nome_shape):
    sp_tree = slide.shapes._spTree
    for shape in slide.shapes:
        if shape.name == nome_shape:
            sp_tree.remove(shape._element)
            return True
    return False


def remover_slide(prs, indice):
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[indice])


def processar_slide(slide, acessos, ordem, topo_inicial, mapa):
    """
    1. Remove os blocos (texto + imagens) dos acessos desativados
    2. Reposiciona verticalmente os blocos que sobraram, de cima para baixo
    """

    # PASSO 1 — remove shapes dos acessos não autorizados
    for acesso, tem_acesso in acessos.items():
        if not tem_acesso:
            remover_shape_por_nome(slide, mapa[acesso]["texto"])
            for img in mapa[acesso]["imagens"]:
                remover_shape_por_nome(slide, img)
            print(f"  ✗ '{acesso}' removido")

    # PASSO 2 — reposiciona os blocos ativos de cima para baixo
    # Filtra só os acessos ativos, mantendo a ordem original do slide
    ativos = [a for a in ordem if acessos.get(a, False)]

    topo_atual = topo_inicial  # começa no topo do primeiro bloco

    for acesso in ativos:
        nome_texto = mapa[acesso]["texto"]
        nomes_imgs = mapa[acesso]["imagens"]

        # Encontra o shape de texto e reposiciona
        for shape in slide.shapes:
            if shape.name == nome_texto:
                altura_bloco = shape.height
                shape.top = topo_atual  # move o bloco de texto para a posição atual

                # Move as imagens para a mesma altura do bloco
                for nome_img in nomes_imgs:
                    for img_shape in slide.shapes:
                        if img_shape.name == nome_img:
                            # Mantém o deslocamento vertical relativo original da imagem
                            # em relação ao bloco de texto (imagem não fica sempre no topo)
                            img_shape.top = topo_atual + (img_shape.top - shape.top) if False else topo_atual
                            break

                topo_atual += altura_bloco + ESPACO_ENTRE_BLOCOS
                break

    if ativos:
        print(f"  ✓ Blocos reposicionados: {' → '.join(ativos)}")


def slide_tem_acesso(acessos):
    return any(acessos.values())


# ─── EXECUÇÃO ─────────────────────────────────────────────────────────────────

prs = Presentation("template_onboarding.pptx")

print(f"\nGerando onboarding para: {NOME}")
print("-" * 40)

# Slide 1 — nome
substituir_texto(prs.slides[0], USUARIO, SENHA, NOME)
print("✓ Slide 1: nome preenchido")

# Slides de acesso
configs = [
    (3, ACESSOS_SLIDE4, ORDEM_SLIDE4, TOPO_INICIAL_SLIDE4, "Slide 4"),
    (4, ACESSOS_SLIDE5, ORDEM_SLIDE5, TOPO_INICIAL_SLIDE5, "Slide 5"),
    (5, ACESSOS_SLIDE6, ORDEM_SLIDE6, TOPO_INICIAL_SLIDE6, "Slide 6"),
]

for idx, acessos, ordem, topo, label in configs:
    print(f"\n{label}:")
    if slide_tem_acesso(acessos):
        processar_slide(prs.slides[idx], acessos, ordem, topo, MAPA_ACESSOS)
        substituir_texto(prs.slides[idx], USUARIO, SENHA, NOME)
        print(f"  ✓ {label} mantido")
    else:
        print(f"  ✗ Sem acessos — slide será removido")

# Remove slides vazios (de trás pra frente para não deslocar índices)
for idx, acessos, _, _, label in reversed(configs):
    if not slide_tem_acesso(acessos):
        remover_slide(prs, idx)
        print(f"\n✗ {label} removido")

# Salva
nome_arquivo = f"Onboarding_{NOME}.pptx"
prs.save(nome_arquivo)

print(f"\n{'='*40}")
print(f"✅ Arquivo gerado: {nome_arquivo}")
print(f"   Total de slides: {len(prs.slides)}")
