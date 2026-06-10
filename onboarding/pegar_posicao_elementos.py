from pptx import Presentation

prs = Presentation("template_onboarding.pptx")

slide = prs.slides[5]  # muda o índice para o slide que quiser inspecionar

for shape in slide.shapes:
    cm = lambda emu: round(emu / 360000, 2)  # converte EMU para cm (mais legível)
    print(f"Nome:   {shape.name}")
    print(f"  top:    {shape.top} EMU  ({cm(shape.top)} cm)")
    print(f"  left:   {shape.left} EMU  ({cm(shape.left)} cm)")
    print(f"  width:  {shape.width} EMU  ({cm(shape.width)} cm)")
    print(f"  height: {shape.height} EMU  ({cm(shape.height)} cm)")
    print()