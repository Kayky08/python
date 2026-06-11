from pptx import Presentation
from pptx.util import Inches
from lxml import etree
import copy

# ─── DADOS DO COLABORADOR ─────────────────────────────────────────────────────
print("=== Gerador de Onboarding ===\n")

NOME    = input("Nome completo: ")
USUARIO = input("Usuário: ")
SENHA   = input("Senha: ")

# ─── DEFINA OS ACESSOS DA PESSOA ──────────────────────────────────────────────
# True  = a pessoa TEM esse acesso (bloco aparece no PDF)
# False = a pessoa NÃO tem (bloco é removido do slide)
#
print("\nDefina os acessos (s = sim / n = não):\n")

def perguntar(acesso):
    resposta = input(f"  {acesso}? (s/n): ").strip().lower()
    return resposta == "s"  # retorna True se digitou "s", False se qualquer outra coisa


# SLIDE 4 — até 3 acessos
ACESSOS_SLIDE4 = {
    "Fluig":     perguntar("Fluig"),
    "Rede":      perguntar("Rede"),
    "Office365": perguntar("Office365"),
}

# SLIDE 5 — até 3 acessos
ACESSOS_SLIDE5 = {
    "SAG":      perguntar("SAG"),
    "Protheus": perguntar("Protheus"),
    "Edata":    perguntar("Edata"),
}

# SLIDE 6 — até 1 acesso
ACESSOS_SLIDE6 = {
    "PowerBI": perguntar("Power BI"),
}

# ─── MAPEAMENTO: qual chave identifica qual shape ──────────────────────────────
# Baseado nos nomes dos shapes inspecionados no template.
# Cada entrada: nome_do_acesso -> nome_do_shape_de_texto no .pptx
MAPA_SHAPES = {
    # Slide 4
    "Fluig":    "object 5",
    "Rede":     "CaixaDeTexto 6",
    "Office365":"CaixaDeTexto 9",
    # Slide 5
    "SAG":      "object 5",
    "Protheus": "CaixaDeTexto 17",
    "Edata":    "CaixaDeTexto 18",
    # Slide 6
    "PowerBI":  "object 5",
}

# Imagens associadas a cada acesso (nome do shape de imagem no .pptx)
# As imagens ficam ao lado de cada bloco de texto — quando o bloco é removido,
# a imagem correspondente também é removida.
MAPA_IMAGENS = {
    "Fluig":    "Imagem 7",    # logo do Fluig no slide 4
    "Office365":"object 6",    # logo do Office365 no slide 4
    "Rede":     "Imagem 5",          # Login de rede não tem imagem própria
    "SAG":      "Imagem 12",
    "Protheus": "Imagem 9",
    "Edata":    "Imagem 14",
    "PowerBI":  "Imagem 5",
}


# ─── FUNÇÕES AUXILIARES ───────────────────────────────────────────────────────

def substituir_texto(slide, usuario, senha):
    """Percorre todos os shapes do slide e substitui os placeholders."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "{{USUARIO}}" in run.text:
                    run.text = run.text.replace("{{USUARIO}}", usuario)
                if "{{SENHA}}" in run.text:
                    run.text = run.text.replace("{{SENHA}}", senha)
                if "{{NOME}}" in run.text:
                    run.text = run.text.replace("{{NOME}}", NOME)


def remover_shape_por_nome(slide, nome_shape):
    """
    Remove um shape do slide pelo seu nome.
    Funciona diretamente no XML — necessário porque python-pptx
    não tem método nativo de remoção de shapes.
    """
    sp_tree = slide.shapes._spTree  # elemento XML raiz do slide (spTree = shape tree)

    for shape in slide.shapes:
        if shape.name == nome_shape:
            sp_tree.remove(shape._element)  # remove o nó XML do elemento
            return True  # encontrou e removeu

    return False  # shape não encontrado (sem erro, apenas avisa)


def remover_acessos_nao_autorizados(slide, mapa_acessos, mapa_shapes, mapa_imagens):
    """
    Para cada acesso marcado como False, remove o bloco de texto
    e a imagem correspondente do slide.
    """
    for acesso, tem_acesso in mapa_acessos.items():
        if not tem_acesso:
            # Remove o bloco de texto do acesso
            nome_shape = mapa_shapes.get(acesso)
            if nome_shape:
                removido = remover_shape_por_nome(slide, nome_shape)
                if removido:
                    print(f"  ✗ Removido: bloco de texto '{acesso}' ({nome_shape})")

            # Remove a imagem associada (se houver)
            nome_imagem = mapa_imagens.get(acesso)
            if nome_imagem:
                removido = remover_shape_por_nome(slide, nome_imagem)
                if removido:
                    print(f"  ✗ Removida: imagem de '{acesso}' ({nome_imagem})")


def slide_tem_algum_acesso(mapa_acessos):
    """Retorna True se pelo menos um acesso do slide está ativo."""
    return any(mapa_acessos.values())


def remover_slide(prs, indice):
    """
    Remove um slide da apresentação pelo índice (0 = primeiro slide).
    Também exige manipulação direta do XML pois python-pptx não tem
    método nativo para isso.
    """
    xml_slides = prs.slides._sldIdLst  # lista de IDs dos slides no XML
    slide_element = xml_slides[indice]  # pega o elemento do slide pelo índice
    xml_slides.remove(slide_element)    # remove da lista de slides


# ─── LÓGICA PRINCIPAL ─────────────────────────────────────────────────────────

prs = Presentation("template_onboarding.pptx")

print(f"\nGerando onboarding para: {NOME}")
print("-" * 40)

# Slide 1 — substitui o nome
substituir_texto(prs.slides[0], USUARIO, SENHA)
print("✓ Slide 1: nome preenchido")

# Slide 4 — acessos: Fluig, Rede, Office365
print("\nSlide 4:")
if slide_tem_algum_acesso(ACESSOS_SLIDE4):
    remover_acessos_nao_autorizados(
        prs.slides[3], ACESSOS_SLIDE4, MAPA_SHAPES, MAPA_IMAGENS
    )
    substituir_texto(prs.slides[3], USUARIO, SENHA)
    print("  ✓ Slide 4 mantido")
else:
    print("  ✗ Nenhum acesso do slide 4 — slide será removido")

# Slide 5 — acessos: SAG, Protheus, Edata
print("\nSlide 5:")
if slide_tem_algum_acesso(ACESSOS_SLIDE5):
    remover_acessos_nao_autorizados(
        prs.slides[4], ACESSOS_SLIDE5, MAPA_SHAPES, MAPA_IMAGENS
    )
    substituir_texto(prs.slides[4], USUARIO, SENHA)
    print("  ✓ Slide 5 mantido")
else:
    print("  ✗ Nenhum acesso do slide 5 — slide será removido")

# Slide 6 — acesso: Power BI
print("\nSlide 6:")
if slide_tem_algum_acesso(ACESSOS_SLIDE6):
    remover_acessos_nao_autorizados(
        prs.slides[5], ACESSOS_SLIDE6, MAPA_SHAPES, MAPA_IMAGENS
    )
    substituir_texto(prs.slides[5], USUARIO, SENHA)
    print("  ✓ Slide 6 mantido")
else:
    print("  ✗ Nenhum acesso do slide 6 — slide será removido")

# ─── REMOVE OS SLIDES VAZIOS ──────────────────────────────────────────────────
# A remoção é feita de trás pra frente para não bagunçar os índices.
# Ex: se removermos o slide 4 primeiro, o slide 5 passa a ser o 4,
# e o índice do slide 6 muda — por isso começamos pelo último.
slides_para_remover = []

if not slide_tem_algum_acesso(ACESSOS_SLIDE6):
    slides_para_remover.append(5)  # índice 5 = slide 6
if not slide_tem_algum_acesso(ACESSOS_SLIDE5):
    slides_para_remover.append(4)
if not slide_tem_algum_acesso(ACESSOS_SLIDE4):
    slides_para_remover.append(3)

# Ordena do maior para o menor índice antes de remover
for idx in sorted(slides_para_remover, reverse=True):
    remover_slide(prs, idx)
    print(f"\n  ✗ Slide {idx+1} removido da apresentação")

# ─── SALVA ────────────────────────────────────────────────────────────────────
nome_arquivo = f"Onboarding_{NOME}.pptx"
prs.save(nome_arquivo)

print(f"\n{'='*40}")
print(f"✅ Arquivo gerado: {nome_arquivo}")
print(f"   Total de slides: {len(prs.slides)}")
