from pptx import Presentation
import copy
import sys

# ─── DADOS DO COLABORADOR ─────────────────────────────────────────────────────
print("=== Gerador de Onboarding ===\n")

#print("Defina os acessos (s = sim / n = não):\n")

def perguntar(acesso):
   return input(f"  {acesso}? (s/n): ").strip().lower() == "s"

# A ordem aqui define a prioridade de preenchimento dos slides
ACESSOS = {
    "Rede":      perguntar("Rede"),
    "Office365": perguntar("Office365"),
    "Fluig":     perguntar("Fluig"),
    "Protheus":  perguntar("Protheus"),
    "Edata":     perguntar("Edata"),
    "SAG":       perguntar("SAG"),
    "PowerBI":   perguntar("Power BI"),
    "WMW":       perguntar("WMW")
}

#---------------Placeholder----------------- 
# ACESSOS = {
#     "Rede":      "s",
#     "Office365": "s",
#     "Fluig":     "s",
#     "Protheus":  "s",
#     "Edata":     "s",
#     "SAG":       "s",
#     "PowerBI":   "s",
# }
#---------------Placeholder----------------- 

#print("\nDefina as informações:\n")

# NOME    = input("Nome completo: ")
# USUARIO = input("Usuário: ")
# SENHA   = input("Senha: ")
# USUARIO_EDATA = ""
# USUARIO_SAG = ""

#---------------Placeholder----------------- 
NOME    = "João da Silva"
USUARIO = "joao.silva"
SENHA   = "Comer@312!"
USUARIO_EDATA = "sjoao"
USUARIO_SAG = "joaosilva"
#---------------Placeholder----------------- 

# if ACESSOS["Edata"]:
#     USUARIO_EDATA = input("Usuário EDATA: ")

# if ACESSOS["SAG"]:
#     USUARIO_SAG = input("Usuário SAG: ")

# ─── DEFINIÇÃO DE CADA ACESSO ─────────────────────────────────────────────────
# texto_shape : nome do shape de texto no template
# imagens     : lista de imagens com nome e offset vertical relativo ao bloco de texto
#               (offset negativo = imagem fica ACIMA do topo do texto)
# slide_origem: índice do slide onde esse acesso vive no template (0 = slide 1)

DEFINICAO_ACESSOS = {
    "Rede": {
        "slide_origem": 3,
        "texto_shape":  "ct_rede",
        "imagens": [
            {"nome": "img_rede", "offset": -482429, "left": 9704070, "w": 1338263, "h": 1765194},  # era -360000
        ],
    },
    "Office365": {
        "slide_origem": 3,
        "texto_shape":  "ct_office",
        "imagens": [
            {"nome": "img_office", "offset": 537266, "left": 9924765, "w": 896873, "h": 807276},
        ],
    },
    "Fluig": {
        "slide_origem": 3,
        "texto_shape":  "ct_fluig",
        "imagens": [
            {"nome": "img_fluig", "offset": 486329, "left": 9704070, "w": 1098803, "h": 516635},
        ],
    },
    "Protheus": {
        "slide_origem": 4,
        "texto_shape":  "ct_protheus",
        "imagens": [
            {"nome": "img_protheus1", "offset": -395287, "left": 9412930,  "w": 714375, "h": 752475},
            {"nome": "img_protheus2", "offset":  428252, "left": 9972893,  "w": 771525, "h": 742950},
            {"nome": "img_protheus3", "offset": -357187, "left": 10550541, "w": 752475, "h": 714375},
        ],
    },
    "Edata": {
        "slide_origem": 4,
        "texto_shape":  "ct_edata",
        "imagens": [
            {"nome": "img_edata1", "offset":  835149, "left": 10045653, "w": 821837, "h": 769657},  # era -66000
            {"nome": "img_edata2", "offset":  -61635, "left": 9513846,  "w": 821837, "h": 731949},
            {"nome": "img_edata3", "offset":  -65836, "left": 10591686, "w": 711330, "h": 731948},
        ],
    },
    "SAG": {
        "slide_origem": 4,
        "texto_shape":  "ct_sag",
        "imagens": [
            {"nome": "img_sag", "offset": 451877, "left": 9940482, "w": 679896, "h": 945942},
        ],
    },
    "PowerBI": {
        "slide_origem": 5,
        "texto_shape":  "ct_bi",
        "imagens": [
            {"nome": "img_bi", "offset": -226036, "left": 9727692, "w": 1473708, "h": 1873896},  # era 0
        ],
    },
    "WMW": {
        "slide_origem": 5,
        "texto_shape":  "ct_wmw",
        "imagens": [],
    },
}

# Capacidade de cada slide e posição inicial dos blocos
CAPACIDADE_POR_SLIDE = 3   # slides 4 e 5 comportam 3 acessos cada
CAPACIDADE_ULTIMO    = 1   # slide 6 comporta 1 acesso
TOPO_INICIAL         = 1496968   # posição vertical do primeiro bloco (em EMU)
ESPACO_ENTRE_BLOCOS  = 200000    # espaço entre blocos consecutivos (em EMU)


# ─── FUNÇÕES AUXILIARES ───────────────────────────────────────────────────────

def substituir_texto(slide, usuario, senha, nome, usuario_edata = "", usuario_sag = ""):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if ACESSOS["Edata"] or ACESSOS["SAG"]:
                    run.text = (run.text
                        .replace("{{USUARIO}}", usuario)
                        .replace("{{SENHA}}", senha)
                        .replace("{{NOME}}", nome)
                        .replace("{{USUARIO_EDATA}}", usuario_edata)
                        .replace("{{USUARIO_SAG}}", usuario_sag))
                else:
                    run.text = (run.text
                        .replace("{{USUARIO}}", usuario)
                        .replace("{{SENHA}}", senha)
                        .replace("{{NOME}}", nome))


def remover_shape_por_nome(slide, nome):
    sp_tree = slide.shapes._spTree
    for shape in slide.shapes:
        if shape.name == nome:
            sp_tree.remove(shape._element)
            return True
    return False


def remover_slide(prs, indice):
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[indice])


def copiar_slide_template(prs, indice_origem):
    """
    Duplica um slide existente e o adiciona ao final da apresentação.
    Necessário quando um grupo de acessos precisa de um slide extra.
    """
    template = prs.slides[indice_origem]
    slide_layout = template.slide_layout

    novo_slide = prs.slides.add_slide(slide_layout)

    # Copia todos os shapes do template para o novo slide
    sp_tree_origem = template.shapes._spTree
    sp_tree_destino = novo_slide.shapes._spTree

    for shape in sp_tree_origem:
        sp_tree_destino.append(copy.deepcopy(shape))

    return novo_slide


def posicionar_bloco(slide, acesso_nome, novo_topo, definicao):
    """
    Move o bloco de texto e todas as suas imagens para novo_topo.
    As imagens mantêm seu offset relativo ao bloco de texto.
    """
    texto_shape_nome = definicao["texto_shape"]

    for shape in slide.shapes:
        if shape.name == texto_shape_nome:
            shape.top = novo_topo

    for img_def in definicao["imagens"]:
        for shape in slide.shapes:
            if shape.name == img_def["nome"]:
                shape.top = novo_topo + img_def["offset"]


def adicionar_acesso_ao_slide(slide, acesso_nome, novo_topo, prs_origem, definicao):
    """
    Copia os shapes de um acesso do slide de origem para o slide de destino,
    já posicionados em novo_topo.
    """
    slide_origem = prs_origem.slides[definicao["slide_origem"]]

    # Copia o bloco de texto
    for shape in slide_origem.shapes:
        if shape.name == definicao["texto_shape"]:
            novo_shape = copy.deepcopy(shape._element)
            slide.shapes._spTree.append(novo_shape)
            # Reposiciona no slide de destino
            for s in slide.shapes:
                if s.name == definicao["texto_shape"]:
                    s.top = novo_topo
            break

    # Copia as imagens
    for img_def in definicao["imagens"]:
        for shape in slide_origem.shapes:
            if shape.name == img_def["nome"]:
                novo_shape = copy.deepcopy(shape._element)
                slide.shapes._spTree.append(novo_shape)
                for s in slide.shapes:
                    if s.name == img_def["nome"]:
                        s.top  = novo_topo + img_def["offset"]
                        s.left = img_def["left"]
                        s.width  = img_def["w"]
                        s.height = img_def["h"]
                break


def limpar_slide_de_acessos(slide):
    """Remove todos os shapes de login de um slide (deixa só o cabeçalho)."""
    shapes_fixos = {"object 2", "object 3"}  # título e barra decorativa
    sp_tree = slide.shapes._spTree
    para_remover = [
        shape._element for shape in slide.shapes
        if shape.name not in shapes_fixos
    ]
    for el in para_remover:
        sp_tree.remove(el)


# ─── LÓGICA PRINCIPAL ─────────────────────────────────────────────────────────

# Filtra só os acessos ativos, mantendo a ordem definida
ativos = [nome for nome, tem in ACESSOS.items() if tem]

print(f"\nGerando onboarding para: {NOME}")
print(f"Acessos ativos ({len(ativos)}): {', '.join(ativos)}")
print("-" * 40)

if not ativos:
    print("Nenhum acesso selecionado.")
    exit()

# Divide os acessos em grupos de até 3, com o último slide aceitando 1
# Exemplo: [Rede, Fluig, Protheus, PowerBI] → [[Rede, Fluig, Protheus], [PowerBI]]
grupos = []
i = 0
while i < len(ativos):
    grupos.append(ativos[i:i + CAPACIDADE_POR_SLIDE])
    i += CAPACIDADE_POR_SLIDE

print(f"Slides de acesso necessários: {len(grupos)}")
for g_idx, grupo in enumerate(grupos):
    print(f"  Slide {4 + g_idx}: {', '.join(grupo)}")

# Abre o template de novo como referência para copiar shapes
prs_ref = Presentation("template_onboarding.pptx")
prs     = Presentation("template_onboarding.pptx")

# Slide 1 — nome
if ACESSOS["Edata"] or ACESSOS["SAG"]:
    substituir_texto(prs.slides[0], USUARIO, SENHA, NOME, USUARIO_EDATA, USUARIO_SAG)
else:
    substituir_texto(prs.slides[0], USUARIO, SENHA, NOME)
    
print("\n✓ Slide 1: nome preenchido")

# Slides de acesso (índices 3, 4, 5 no template = slides 4, 5, 6)
slides_acesso_template = [3, 4, 5]  # índices no prs original

# Quantos slides de acesso serão usados
slides_necessarios = len(grupos)
slides_disponiveis = len(slides_acesso_template)  # sempre 3 no template

# Processa cada grupo de acessos
for g_idx, grupo in enumerate(grupos):
    slide_idx = slides_acesso_template[g_idx]
    slide = prs.slides[slide_idx]

    print(f"\nSlide {4 + g_idx} — {', '.join(grupo)}:")

    # Remove todos os blocos de acesso do slide (começa limpo)
    limpar_slide_de_acessos(slide)

    # Adiciona e posiciona cada acesso do grupo
    topo_atual = TOPO_INICIAL
    for acesso in grupo:
        definicao = DEFINICAO_ACESSOS[acesso]
        adicionar_acesso_ao_slide(slide, acesso, topo_atual, prs_ref, definicao)
        # Avança o topo para o próximo bloco
        for shape in slide.shapes:
            if shape.name == definicao["texto_shape"]:
                topo_atual += shape.height + ESPACO_ENTRE_BLOCOS
                break
        print(f"  ✓ '{acesso}' posicionado")

    substituir_texto(slide, USUARIO, SENHA, NOME, USUARIO_EDATA, USUARIO_SAG)

# Remove slides de acesso que sobraram (de trás pra frente)
slides_para_remover = slides_acesso_template[slides_necessarios:]
for idx in sorted(slides_para_remover, reverse=True):
    remover_slide(prs, idx)
    print(f"\n✗ Slide {idx + 1} removido (não necessário)")

# Salva
nome_arquivo = f"Onboarding - {NOME}.pptx"
prs.save(nome_arquivo)

print(f"\n{'='*40}")
print(f"✅ Arquivo gerado: {nome_arquivo}")
print(f"   Total de slides: {len(prs.slides)}")
