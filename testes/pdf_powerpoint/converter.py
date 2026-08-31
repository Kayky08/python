import io
import pymupdf  # Nome de importação da biblioteca PyMuPDF
from pptx import Presentation
from pptx.util import Inches

def pdf_para_powerpoint(caminho_pdf, caminho_pptx):
    # Abre o arquivo PDF
    doc = pymupdf.open(caminho_pdf)
    prs = Presentation()
    
    # Remove o slide em branco inicial padrão
    if prs.slides:
        prs.slides._sldIdLst.clear()

    for num_pagina in range(len(doc)):
        pagina = doc.load_page(num_pagina)
        
        # Renderiza a página como imagem (DPI 150 garante boa qualidade sem pesar)
        pix = pagina.get_pixmap(dpi=150)
        dados_imagem = pix.tobytes("png")
        
        # Ajusta o tamanho do slide para o tamanho exato da página do PDF (72 pontos por polegada)
        prs.slide_width = Inches(pagina.rect.width / 72)
        prs.slide_height = Inches(pagina.rect.height / 72)
        
        # Cria um slide totalmente em branco
        layout_branco = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout_branco)
        
        # Insere a imagem preenchendo o slide inteiro
        fluxo_imagem = io.BytesIO(dados_imagem)
        slide.shapes.add_picture(fluxo_imagem, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Salva o arquivo final
    prs.save(caminho_pptx)
    print(f"\nConversão concluída com sucesso! Arquivo salvo em: {caminho_pptx}")

# Defina os nomes dos seus arquivos aqui
pdf_para_powerpoint("TA2 - Fluxo de reprocessamento de material.pdf", "TA2 - Fluxo de reprocessamento de material.pptx")