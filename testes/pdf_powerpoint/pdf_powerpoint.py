import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path


def pdf_para_pptx(pdf_path, pptx_path):
    # Abre o PDF
    pdf = fitz.open(pdf_path)

    # Cria uma apresentação
    prs = Presentation()

    # Remove o slide inicial padrão
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        rId = slide.part.relate_to(
            slide.background.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
        )

    for pagina in pdf:
        # Renderiza a página como imagem
        pix = pagina.get_pixmap(matrix=fitz.Matrix(2, 2))

        imagem = "pagina_temp.png"
        pix.save(imagem)

        # Adiciona um slide em branco
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Tamanho da página do PDF
        largura_pdf = pagina.rect.width
        altura_pdf = pagina.rect.height

        # Ajusta o tamanho do slide para manter a proporção
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(
            13.333 * altura_pdf / largura_pdf
        )

        # Adiciona a imagem ocupando todo o slide
        slide.shapes.add_picture(
            imagem,
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        # Apaga a imagem temporária
        Path(imagem).unlink()

    # Salva o PowerPoint
    prs.save(pptx_path)
    pdf.close()


# -----------------------------
# USO
# -----------------------------

pdf_para_pptx(
    "meu_arquivo.pdf",
    "meu_arquivo.pptx"
)

print("PDF convertido com sucesso!")